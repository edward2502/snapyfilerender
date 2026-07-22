"""
converter.py - Core conversion logic used by the API.

Word <-> PDF: LibreOffice does docx->pdf; pdf2docx does pdf->docx.
Excel <-> PDF: LibreOffice does xlsx->pdf (with auto column-widening and
               thin cell borders added first, so the exported PDF has a
               real grid instead of ambiguous plain text); pdfplumber
               extracts tables back out of a PDF into an .xlsx workbook.
"""

import os
import shutil
import subprocess
import tempfile


def find_soffice():
    for name in ("soffice", "libreoffice"):
        path = shutil.which(name)
        if path:
            return path
    candidates = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
    ]
    for c in candidates:
        if os.path.exists(c):
            return c
    return None


def _office_to_pdf(input_path: str, output_path: str) -> str:
    soffice = find_soffice()
    if not soffice:
        raise RuntimeError("LibreOffice (soffice) not found on this server.")

    with tempfile.TemporaryDirectory() as tmpdir:
        cmd = [
            soffice, "--headless", "--norestore",
            "--convert-to", "pdf",
            "--outdir", tmpdir,
            input_path,
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

        produced = os.path.join(
            tmpdir, os.path.splitext(os.path.basename(input_path))[0] + ".pdf"
        )
        if not os.path.exists(produced):
            raise RuntimeError("Conversion did not produce a PDF.")
        shutil.move(produced, output_path)

    return output_path


def word_to_pdf(input_path: str, output_path: str) -> str:
    return _office_to_pdf(input_path, output_path)


def pdf_to_word(input_path: str, output_path: str) -> str:
    from pdf2docx import Converter

    cv = Converter(input_path)
    try:
        cv.convert(output_path)
    finally:
        cv.close()

    return output_path


def _autofit_and_border_xlsx(input_path, working_copy_path):
    """Widen columns that would otherwise get clipped on PDF export, and
    add thin borders so the exported PDF has a real grid - which makes
    pdf_to_excel's extraction exact instead of guesswork."""
    import openpyxl
    from openpyxl.styles import Border, Side

    thin = Side(style="thin")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    wb = openpyxl.load_workbook(input_path)
    for ws in wb.worksheets:
        for col_cells in ws.columns:
            longest = max((len(str(c.value)) for c in col_cells if c.value is not None), default=0)
            if longest == 0:
                continue
            col_letter = col_cells[0].column_letter
            target_width = max(longest + 2, 8)
            current = ws.column_dimensions[col_letter].width
            if current is None or current < target_width:
                ws.column_dimensions[col_letter].width = target_width

        if ws.dimensions and ws.dimensions != "A1:A1":
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        cell.border = border

    wb.save(working_copy_path)
    return working_copy_path


def excel_to_pdf(input_path: str, output_path: str) -> str:
    with tempfile.TemporaryDirectory() as tmpdir:
        widened = os.path.join(tmpdir, os.path.basename(input_path))
        try:
            _autofit_and_border_xlsx(input_path, widened)
            source = widened
        except Exception:
            source = input_path
        return _office_to_pdf(source, output_path)


def _extract_table_from_words(words, row_tolerance=4):
    """Fallback table reconstruction for PDFs with no visible ruling lines."""
    if not words:
        return []

    rows = {}
    for w in words:
        key = round(w["top"] / row_tolerance)
        rows.setdefault(key, []).append(w)

    row_keys = sorted(rows.keys())
    ordered_rows = [sorted(rows[k], key=lambda w: w["x0"]) for k in row_keys]
    if not ordered_rows:
        return []

    ncols = len(ordered_rows[0])
    if ncols == 0:
        return []

    def row_to_cells(row_words):
        cells = [dict(w) for w in row_words]
        while len(cells) > ncols:
            gaps = [cells[i + 1]["x0"] - cells[i]["x1"] for i in range(len(cells) - 1)]
            i = gaps.index(min(gaps))
            merged_text = cells[i]["text"] + " " + cells[i + 1]["text"]
            cells[i] = {"text": merged_text, "x0": cells[i]["x0"], "x1": cells[i + 1]["x1"]}
            del cells[i + 1]
        while len(cells) < ncols:
            cells.append({"text": "", "x0": None, "x1": None})
        return cells

    rows_of_cells = [row_to_cells(r) for r in ordered_rows]

    if len(rows_of_cells) < 2 or ncols < 2:
        return []

    for col in range(ncols):
        x0s = [row[col]["x0"] for row in rows_of_cells if row[col]["x0"] is not None]
        if len(x0s) < 2:
            continue
        if max(x0s) - min(x0s) > 12:
            return []

    return [[c["text"] for c in row] for row in rows_of_cells]


def pdf_to_excel(input_path: str, output_path: str) -> str:
    import pdfplumber
    import pandas as pd

    page_frames = {}
    with pdfplumber.open(input_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            tables = page.extract_tables()
            frames = []

            if tables:
                for table in tables:
                    if not table or len(table) < 1:
                        continue
                    header, *rows = table
                    frames.append(pd.DataFrame(rows, columns=header))
            else:
                words = page.extract_words()
                table = _extract_table_from_words(words)
                if len(table) > 1:
                    header, *rows = table
                    frames.append(pd.DataFrame(rows, columns=header))

            if frames:
                page_frames[page_num] = frames

    if not page_frames:
        raise RuntimeError(
            "No tables were found in this PDF. Excel conversion only works "
            "on PDFs that actually contain tabular data."
        )

    with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
        for page_num, frames in page_frames.items():
            sheet_name = f"Page{page_num}"[:31]
            start_row = 0
            for df in frames:
                df.to_excel(writer, sheet_name=sheet_name, index=False, startrow=start_row)
                start_row += len(df) + 2

    return output_path


def merge_pdfs(input_paths: list, output_path: str) -> str:
    """Combine multiple PDFs into one, in the given order."""
    from pypdf import PdfWriter

    writer = PdfWriter()
    for path in input_paths:
        writer.append(path)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def split_pdf(input_path: str, output_path: str, start_page: int, end_page: int) -> str:
    """Extract a page range (1-indexed, inclusive) into a new PDF."""
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(input_path)
    total_pages = len(reader.pages)

    if start_page < 1 or end_page < 1:
        raise ValueError("Page numbers must be 1 or greater.")
    if start_page > total_pages:
        raise ValueError(f"This PDF only has {total_pages} page(s).")
    if end_page > total_pages:
        end_page = total_pages
    if start_page > end_page:
        raise ValueError("Start page must be less than or equal to end page.")

    writer = PdfWriter()
    for i in range(start_page - 1, end_page):
        writer.add_page(reader.pages[i])

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def compress_pdf(input_path: str, output_path: str, image_quality: int = 60, max_dimension: int = 1600) -> str:
    """Shrink a PDF's file size by re-encoding its embedded images at a
    lower quality/resolution and compressing internal object streams.
    Text and vector content are left untouched - only raster images (the
    usual cause of PDF bloat) are recompressed."""
    import io
    import pikepdf
    from PIL import Image

    pdf = pikepdf.open(input_path)

    for page in pdf.pages:
        if "/Resources" not in page or "/XObject" not in page.Resources:
            continue
        xobjects = page.Resources.XObject
        for name in list(xobjects.keys()):
            obj = xobjects[name]
            if obj.get("/Subtype") != pikepdf.Name("/Image"):
                continue
            try:
                pdf_image = pikepdf.PdfImage(obj)
                pil_image = pdf_image.as_pil_image()
            except Exception:
                continue  # unsupported image encoding - leave it as-is

            # Downscale oversized images
            w, h = pil_image.size
            if max(w, h) > max_dimension:
                scale = max_dimension / max(w, h)
                pil_image = pil_image.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

            if pil_image.mode not in ("RGB", "L"):
                pil_image = pil_image.convert("RGB")

            buf = io.BytesIO()
            pil_image.save(buf, format="JPEG", quality=image_quality, optimize=True)
            buf.seek(0)

            try:
                new_image = pikepdf.PdfImage.make_image(pdf, buf.read(), pil_image.size)
            except AttributeError:
                # Fallback for pikepdf versions without make_image helper:
                # replace the stream data directly with a DCTDecode (JPEG) filter.
                obj.write(buf.getvalue(), filter=pikepdf.Name("/DCTDecode"))
                obj.Width = pil_image.width
                obj.Height = pil_image.height
                obj.ColorSpace = pikepdf.Name("/DeviceRGB") if pil_image.mode == "RGB" else pikepdf.Name("/DeviceGray")
                obj.BitsPerComponent = 8
                continue

            xobjects[name] = new_image

    pdf.save(output_path, compress_streams=True, object_stream_mode=pikepdf.ObjectStreamMode.generate)
    pdf.close()
    return output_path


def unlock_pdf(input_path: str, output_path: str, password: str) -> str:
    """Remove password protection from a PDF, given the correct password."""
    from pypdf import PdfReader, PdfWriter
    from pypdf.errors import FileNotDecryptedError

    reader = PdfReader(input_path)
    if reader.is_encrypted:
        result = reader.decrypt(password)
        if result == 0:
            raise ValueError("Incorrect password.")

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def protect_pdf(input_path: str, output_path: str, password: str) -> str:
    """Add password protection to a PDF."""
    from pypdf import PdfReader, PdfWriter

    if not password:
        raise ValueError("Password cannot be empty.")

    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    writer.encrypt(user_password=password, algorithm="AES-256")

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def watermark_pdf(input_path: str, output_path: str, text: str, opacity: float = 0.3) -> str:
    """Stamp a diagonal text watermark across every page."""
    import io
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    from reportlab.lib.colors import Color

    reader = PdfReader(input_path)
    writer = PdfWriter()

    for page in reader.pages:
        page_width = float(page.mediabox.width)
        page_height = float(page.mediabox.height)

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=(page_width, page_height))
        c.saveState()
        c.translate(page_width / 2, page_height / 2)
        c.rotate(45)
        c.setFillColor(Color(0.5, 0.5, 0.5, alpha=opacity))
        font_size = max(24, min(page_width, page_height) / 10)
        c.setFont("Helvetica-Bold", font_size)
        c.drawCentredString(0, 0, text)
        c.restoreState()
        c.save()
        buf.seek(0)

        watermark_reader = PdfReader(buf)
        page.merge_page(watermark_reader.pages[0])
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def generate_qr_code(data: str, output_path: str, box_size: int = 10, fill_color: str = "black", back_color: str = "white") -> str:
    """Generate a QR code PNG encoding the given text/URL."""
    import qrcode

    if not data or not data.strip():
        raise ValueError("QR code content cannot be empty.")

    qr = qrcode.QRCode(
        version=None,  # auto-size to fit the data
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=box_size,
        border=4,
    )
    qr.add_data(data)
    qr.make(fit=True)

    img = qr.make_image(fill_color=fill_color, back_color=back_color)
    img.save(output_path)
    return output_path


def extract_pdf_text(input_path: str, max_chars: int = 100_000) -> str:
    """Pull all text out of a PDF for use as AI context. Capped at
    max_chars so a huge PDF can't blow past the model's context window
    or balloon token usage on the free tier."""
    import pdfplumber

    text_parts = []
    total_len = 0
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            text_parts.append(page_text)
            total_len += len(page_text)
            if total_len >= max_chars:
                break

    full_text = "\n\n".join(text_parts)
    if len(full_text) > max_chars:
        full_text = full_text[:max_chars] + "\n\n[... document truncated for length ...]"

    if not full_text.strip():
        raise ValueError(
            "No text could be extracted from this PDF. It may be a scanned "
            "image without a text layer - Chat with PDF only works on PDFs "
            "with selectable text."
        )
    return full_text


def chat_with_pdf(document_text: str, question: str, api_key: str) -> str:
    """Answer a question about a PDF's contents using Google's Gemini API.
    document_text should already be extracted (see extract_pdf_text) so
    that a multi-turn conversation doesn't need to re-parse the PDF on
    every question."""
    from google import genai

    if not question or not question.strip():
        raise ValueError("Question cannot be empty.")
    if not api_key:
        raise RuntimeError(
            "No Gemini API key configured on the server (GEMINI_API_KEY env var missing)."
        )

    client = genai.Client(api_key=api_key)

    prompt = (
        "You are answering questions about a specific document. Only use "
        "the document text below to answer - if the answer isn't in the "
        "document, say so clearly rather than guessing.\n\n"
        "--- DOCUMENT START ---\n"
        f"{document_text}\n"
        "--- DOCUMENT END ---\n\n"
        f"Question: {question}"
    )

    response = client.models.generate_content(
        model="gemini-2.0-flash",
        contents=prompt,
    )
    return response.text


def images_to_pdf(input_paths: list, output_path: str) -> str:
    """Combine one or more images (jpg/png/etc.) into a single PDF, one
    image per page, in the given order."""
    from PIL import Image

    if not input_paths:
        raise ValueError("At least one image is required.")

    images = []
    for path in input_paths:
        img = Image.open(path)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        images.append(img)

    first, rest = images[0], images[1:]
    first.save(output_path, save_all=True, append_images=rest)
    return output_path


def pdf_to_images(input_path: str, output_dir: str, dpi: int = 150, fmt: str = "png") -> list:
    """Render every page of a PDF to a separate image file. Returns the
    list of output file paths, one per page, in page order."""
    import fitz
    import os

    doc = fitz.open(input_path)
    if doc.page_count == 0:
        raise ValueError("This PDF has no pages.")

    output_paths = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=dpi)
        out_path = os.path.join(output_dir, f"page_{i + 1}.{fmt}")
        pix.save(out_path)
        output_paths.append(out_path)

    doc.close()
    return output_paths


def compress_image(input_path: str, output_path: str, quality: int = 60, max_dimension: int = None) -> str:
    """Shrink an image's file size by re-encoding as JPEG at a lower
    quality, and optionally downscaling if it exceeds max_dimension."""
    from PIL import Image

    img = Image.open(input_path)
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")

    if max_dimension:
        w, h = img.size
        if max(w, h) > max_dimension:
            scale = max_dimension / max(w, h)
            img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

    img.save(output_path, format="JPEG", quality=quality, optimize=True)
    return output_path


def upscale_image(input_path: str, output_path: str, scale_factor: float = 2.0) -> str:
    """Increase an image's resolution using high-quality Lanczos
    resampling plus a mild sharpening pass. Note: this is classic
    upscaling, not AI super-resolution - it enlarges cleanly without
    introducing artifacts, but it cannot invent detail that was never in
    the original image."""
    from PIL import Image, ImageFilter

    if scale_factor <= 1.0:
        raise ValueError("Scale factor must be greater than 1.0 to upscale.")
    if scale_factor > 4.0:
        raise ValueError("Scale factor cannot exceed 4x (to keep processing time and file size reasonable).")

    img = Image.open(input_path)
    if img.mode not in ("RGB", "L", "RGBA"):
        img = img.convert("RGB")

    w, h = img.size
    new_size = (int(w * scale_factor), int(h * scale_factor))

    # Cap absolute output size to avoid runaway memory use on huge inputs
    max_total_pixels = 40_000_000  # ~40MP
    if new_size[0] * new_size[1] > max_total_pixels:
        raise ValueError(
            f"Requested output size ({new_size[0]}x{new_size[1]}) is too large. "
            "Try a smaller scale factor or a smaller source image."
        )

    upscaled = img.resize(new_size, Image.LANCZOS)
    sharpened = upscaled.filter(ImageFilter.UnsharpMask(radius=2, percent=60, threshold=2))

    save_kwargs = {}
    fmt = (img.format or "PNG").upper()
    if fmt in ("JPEG", "JPG"):
        save_kwargs = {"quality": 95, "optimize": True}
        if sharpened.mode == "RGBA":
            sharpened = sharpened.convert("RGB")

    sharpened.save(output_path, **save_kwargs)
    return output_path


def pptx_to_pdf(input_path: str, output_path: str) -> str:
    """Convert a PowerPoint file to PDF using the same LibreOffice engine
    already used for Word and Excel."""
    return _office_to_pdf(input_path, output_path)


def pdf_to_pptx(input_path: str, output_path: str, dpi: int = 150) -> str:
    """Convert a PDF into a PowerPoint file by rendering each page as a
    high-resolution image and placing one image per slide, sized to fill
    the slide exactly.

    Note: this produces a visually faithful presentation that opens and
    presents correctly in PowerPoint/Google Slides, but the text on each
    slide is NOT selectable or editable - each slide is really a picture.
    There is no reliable way to reconstruct genuinely editable text boxes
    and shapes from an arbitrary PDF's layout; every converter that claims
    to do this either produces exactly this kind of image-based result
    under the hood, or badly mangles complex layouts trying to guess at
    "real" shapes."""
    import fitz
    from pptx import Presentation
    from pptx.util import Emu

    doc = fitz.open(input_path)
    if doc.page_count == 0:
        raise ValueError("This PDF has no pages.")

    prs = Presentation()
    # EMU (English Metric Units) - the unit python-pptx uses for sizing
    EMU_PER_INCH = 914400

    with tempfile.TemporaryDirectory() as tmpdir:
        blank_layout = prs.slide_layouts[6]  # fully blank layout

        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=dpi)
            img_path = os.path.join(tmpdir, f"page_{i}.png")
            pix.save(img_path)

            page_width_in = page.rect.width / 72  # PDF points -> inches
            page_height_in = page.rect.height / 72

            if i == 0:
                prs.slide_width = Emu(int(page_width_in * EMU_PER_INCH))
                prs.slide_height = Emu(int(page_height_in * EMU_PER_INCH))

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                img_path, 0, 0,
                width=prs.slide_width, height=prs.slide_height,
            )

        prs.save(output_path)

    doc.close()
    return output_path


def powerpoint_to_pdf(input_path: str, output_path: str) -> str:
    """Convert a .pptx to .pdf using the same LibreOffice engine as
    Word/Excel conversion."""
    return _office_to_pdf(input_path, output_path)


def pdf_to_powerpoint(input_path: str, output_path: str, dpi: int = 150) -> str:
    """Convert a PDF into a .pptx, one slide per page.

    Honesty note: LibreOffice has no PDF->PPTX export filter at all, and
    there is no reliable way to reconstruct fully editable text boxes
    from an arbitrary PDF's layout. This renders each page as a
    high-quality image and places it as a full-slide picture - visually
    identical to the original, but the text on each slide is not
    editable. This is the same practical tradeoff most PDF-to-PowerPoint
    tools make under the hood for anything beyond simple text PDFs.
    """
    import tempfile
    from pptx import Presentation
    from pptx.util import Emu
    from PIL import Image

    with tempfile.TemporaryDirectory() as tmpdir:
        image_paths = pdf_to_images(input_path, tmpdir, dpi=dpi)

        prs = Presentation()
        prs.slide_width = Emu(12192000)   # 16:9 widescreen
        prs.slide_height = Emu(6858000)
        blank_layout = prs.slide_layouts[6]

        for img_path in image_paths:
            slide = prs.slides.add_slide(blank_layout)
            img = Image.open(img_path)
            img_w, img_h = img.size
            aspect = img_w / img_h
            slide_aspect = prs.slide_width / prs.slide_height

            if aspect > slide_aspect:
                w = prs.slide_width
                h = int(w / aspect)
            else:
                h = prs.slide_height
                w = int(h * aspect)

            left = (prs.slide_width - w) // 2
            top = (prs.slide_height - h) // 2
            slide.shapes.add_picture(img_path, left, top, width=w, height=h)

        prs.save(output_path)

    return output_path


# ---------------------------------------------------------------------------
# Newer PDF tools: remove/organize/extract pages, scan (OCR), repair,
# html->pdf, page numbers, crop, translate, pdf->markdown, redact.
# ---------------------------------------------------------------------------

def remove_pages(input_path: str, output_path: str, pages_to_remove: list) -> str:
    """Delete the given pages (1-indexed) and keep the rest, in order."""
    from pypdf import PdfReader, PdfWriter

    if not pages_to_remove:
        raise ValueError("Please specify at least one page to remove.")

    reader = PdfReader(input_path)
    total_pages = len(reader.pages)
    remove_set = set()
    for p in pages_to_remove:
        if p < 1 or p > total_pages:
            raise ValueError(f"Page {p} is out of range - this PDF has {total_pages} page(s).")
        remove_set.add(p)

    if len(remove_set) == total_pages:
        raise ValueError("Cannot remove every page from the PDF.")

    writer = PdfWriter()
    for i, page in enumerate(reader.pages, start=1):
        if i not in remove_set:
            writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def organize_pdf(input_path: str, output_path: str, page_order: list) -> str:
    """Rebuild a PDF using an explicit new page order (1-indexed). Pages
    can be reordered, dropped (by omitting them), or repeated - this
    powers a drag-and-drop 'Organize PDF' UI where the user rearranges
    page thumbnails and/or deletes some."""
    from pypdf import PdfReader, PdfWriter

    if not page_order:
        raise ValueError("Please provide the new page order.")

    reader = PdfReader(input_path)
    total_pages = len(reader.pages)

    writer = PdfWriter()
    for p in page_order:
        if p < 1 or p > total_pages:
            raise ValueError(f"Page {p} is out of range - this PDF has {total_pages} page(s).")
        writer.add_page(reader.pages[p - 1])

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def extract_pages(input_path: str, output_path: str, pages: list) -> str:
    """Pull out an arbitrary set of pages (1-indexed, any order,
    duplicates allowed) into a new PDF."""
    from pypdf import PdfReader, PdfWriter

    if not pages:
        raise ValueError("Please specify at least one page to extract.")

    reader = PdfReader(input_path)
    total_pages = len(reader.pages)

    writer = PdfWriter()
    for p in pages:
        if p < 1 or p > total_pages:
            raise ValueError(f"Page {p} is out of range - this PDF has {total_pages} page(s).")
        writer.add_page(reader.pages[p - 1])

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def repair_pdf(input_path: str, output_path: str) -> str:
    """Attempt to fix a corrupted or malformed PDF by re-parsing it with
    pikepdf (built on qpdf, which repairs broken cross-reference tables,
    broken object streams, and other structural issues) and writing out
    a clean copy."""
    import pikepdf

    try:
        pdf = pikepdf.open(input_path)
    except Exception as e:
        raise ValueError(f"This PDF is too damaged to repair automatically: {e}")

    pdf.save(output_path)
    pdf.close()
    return output_path


def html_to_pdf(input_path: str, output_path: str) -> str:
    """Convert an HTML file to PDF using the same LibreOffice engine used
    for Word/Excel/PowerPoint conversion."""
    return _office_to_pdf(input_path, output_path)


def add_page_numbers(
    input_path: str,
    output_path: str,
    position: str = "bottom-center",
    start_at: int = 1,
    font_size: int = 11,
) -> str:
    """Stamp a page number onto every page of a PDF."""
    import io
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas

    valid_positions = {
        "bottom-center", "bottom-left", "bottom-right",
        "top-center", "top-left", "top-right",
    }
    if position not in valid_positions:
        raise ValueError(f"Position must be one of: {', '.join(sorted(valid_positions))}")

    reader = PdfReader(input_path)
    writer = PdfWriter()

    for i, page in enumerate(reader.pages):
        page_width = float(page.mediabox.width)
        page_height = float(page.mediabox.height)
        number = start_at + i

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=(page_width, page_height))
        c.setFont("Helvetica", font_size)
        margin = 24
        label = str(number)

        y = margin if "bottom" in position else page_height - margin

        if "center" in position:
            c.drawCentredString(page_width / 2, y, label)
        elif "left" in position:
            c.drawString(margin, y, label)
        else:
            c.drawRightString(page_width - margin, y, label)

        c.save()
        buf.seek(0)

        overlay_reader = PdfReader(buf)
        page.merge_page(overlay_reader.pages[0])
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def crop_pdf(input_path: str, output_path: str, left: float = 0, top: float = 0, right: float = 0, bottom: float = 0) -> str:
    """Trim a fixed margin (in points - 72pt = 1 inch) off each side of
    every page."""
    from pypdf import PdfReader, PdfWriter
    from pypdf.generic import RectangleObject

    if any(m < 0 for m in (left, top, right, bottom)):
        raise ValueError("Margins cannot be negative.")

    reader = PdfReader(input_path)
    writer = PdfWriter()

    for page in reader.pages:
        box = page.mediabox
        new_left = float(box.left) + left
        new_bottom = float(box.bottom) + bottom
        new_right = float(box.right) - right
        new_top = float(box.top) - top

        if new_right <= new_left or new_top <= new_bottom:
            raise ValueError("Crop margins are too large for this page size.")

        rect = RectangleObject((new_left, new_bottom, new_right, new_top))
        page.mediabox = rect
        page.cropbox = rect
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def translate_pdf(input_path: str, output_path: str, target_language: str, api_key: str) -> str:
    """Extract a PDF's text, translate it into the target language using
    Gemini, and lay the translation out as a new PDF."""
    from google import genai
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from xml.sax.saxutils import escape

    if not target_language or not target_language.strip():
        raise ValueError("Target language is required.")
    if not api_key:
        raise RuntimeError("No Gemini API key configured on the server (GEMINI_API_KEY env var missing).")

    document_text = extract_pdf_text(input_path, max_chars=100_000)

    client = genai.Client(api_key=api_key)
    prompt = (
        f"Translate the following document text into {target_language}. "
        "Preserve paragraph breaks. Only output the translated text, "
        "nothing else - no notes, no commentary, no markdown.\n\n"
        f"{document_text}"
    )
    response = client.models.generate_content(model="gemini-2.0-flash", contents=prompt)
    translated_text = (response.text or "").strip()

    if not translated_text:
        raise RuntimeError("Translation failed - no text was returned.")

    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(
        output_path, pagesize=A4,
        leftMargin=0.9 * inch, rightMargin=0.9 * inch,
        topMargin=0.9 * inch, bottomMargin=0.9 * inch,
    )
    story = []
    for para in translated_text.split("\n\n"):
        para = para.strip()
        if not para:
            continue
        safe = escape(para).replace("\n", "<br/>")
        story.append(Paragraph(safe, styles["Normal"]))
        story.append(Spacer(1, 10))

    if not story:
        raise RuntimeError("Nothing to render after translation.")

    doc.build(story)
    return output_path


def pdf_to_markdown(input_path: str, output_path: str) -> str:
    """Convert a PDF's text content into a Markdown file, one section per
    page. This is a straightforward text extraction, not a full layout
    reconstruction - tables/images are not preserved as Markdown tables
    or embedded images."""
    import pdfplumber

    lines = []
    with pdfplumber.open(input_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            lines.append(f"## Page {i}\n")
            lines.append(text if text else "*[No extractable text on this page]*")
            lines.append("\n")

    markdown = "\n".join(lines).strip() + "\n"
    if not markdown.strip():
        raise ValueError("No extractable text found in this PDF.")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(markdown)
    return output_path


def redact_pdf(input_path: str, output_path: str, terms: list) -> str:
    """Permanently black out every occurrence of the given search terms.
    Uses PyMuPDF's redaction annotations, which remove the underlying
    text and image content beneath the box - not just draw a black
    rectangle on top of it - so the redacted text cannot be recovered
    by copy-pasting or re-extracting."""
    import fitz

    clean_terms = [t.strip() for t in terms if t and t.strip()]
    if not clean_terms:
        raise ValueError("Please provide at least one word or phrase to redact.")

    doc = fitz.open(input_path)
    if doc.page_count == 0:
        raise ValueError("This PDF has no pages.")

    match_count = 0
    for page in doc:
        for term in clean_terms:
            instances = page.search_for(term)
            for inst in instances:
                page.add_redact_annot(inst, fill=(0, 0, 0))
                match_count += 1
        page.apply_redactions()

    if match_count == 0:
        doc.close()
        raise ValueError("None of the given terms were found in this PDF.")

    doc.save(output_path)
    doc.close()
    return output_path


def scan_pdf(input_path: str, output_path: str, language: str = "eng", dpi: int = 300) -> str:
    """OCR a scanned/image-only PDF and add an invisible, searchable text
    layer on top of the original page images - the pages still look
    identical, but the text can now be selected, copy-pasted, and
    searched. Requires the Tesseract OCR engine to be installed on the
    server (apt package tesseract-ocr)."""
    import io
    import fitz
    from PIL import Image
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas

    try:
        import pytesseract
        from pytesseract import Output
    except ImportError:
        raise RuntimeError("OCR support (pytesseract) is not installed on this server.")

    doc = fitz.open(input_path)
    if doc.page_count == 0:
        raise ValueError("This PDF has no pages.")

    reader = PdfReader(input_path)
    writer = PdfWriter()

    for i, fitz_page in enumerate(doc):
        pix = fitz_page.get_pixmap(dpi=dpi)
        pil_img = Image.open(io.BytesIO(pix.tobytes("png")))

        try:
            ocr_data = pytesseract.image_to_data(pil_img, lang=language, output_type=Output.DICT)
        except pytesseract.TesseractNotFoundError:
            raise RuntimeError("OCR engine (Tesseract) is not installed on this server.")

        pdf_page = reader.pages[i]
        page_width = float(pdf_page.mediabox.width)
        page_height = float(pdf_page.mediabox.height)
        scale_x = page_width / pix.width
        scale_y = page_height / pix.height

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=(page_width, page_height))
        n_boxes = len(ocr_data["text"])
        for j in range(n_boxes):
            word = ocr_data["text"][j].strip()
            if not word:
                continue
            x, y, w, h = (ocr_data["left"][j], ocr_data["top"][j],
                          ocr_data["width"][j], ocr_data["height"][j])
            pdf_x = x * scale_x
            pdf_y = page_height - (y + h) * scale_y  # flip y-axis (image top-down -> PDF bottom-up)
            font_size = max(4, h * scale_y)

            text_obj = c.beginText(pdf_x, pdf_y)
            text_obj.setFont("Helvetica", font_size)
            text_obj.setTextRenderMode(3)  # invisible - present for search/copy, not drawn
            text_obj.textOut(word)
            c.drawText(text_obj)

        c.save()
        buf.seek(0)

        overlay_reader = PdfReader(buf)
        pdf_page.merge_page(overlay_reader.pages[0])
        writer.add_page(pdf_page)

    doc.close()
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path
